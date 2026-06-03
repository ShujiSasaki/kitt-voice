"""
BTC danjer判断AI プレゼン v3 — Project danjer-GAIA
================================================================
3者会議 v3 (Round 0-16) の合意を完全反映:
- Phase 1: 戦略骨格21項目+Appendix A-E+リスク表R1-R33
- Phase 2: 動線「目的→判断→防衛→実装→決断」、 9枚構成
- Phase 3 (本ファイル): Gemini Round 16 ワイヤフレーム + GPT必須修正指示

カラー: NAVY/GOLD/RED/WHITE/MUTED/GRAY (Geminiパレット固定)
タイポ: Hiragino Sans 60/32/24/14pt
余白: 外周80px
"""
import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Circle, Rectangle, Polygon, Wedge
from matplotlib.lines import Line2D
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

plt.rcParams['font.family'] = ['Hiragino Sans', 'Yu Gothic', 'IPAGothic', 'sans-serif']
plt.rcParams['font.size'] = 12
plt.rcParams['axes.unicode_minus'] = False

# Geminiワイヤフレーム カラーパレット
BG_NAVY      = '#0B132B'
TEXT_WHITE   = '#FFFFFF'
TEXT_MUTED   = '#AFB8C1'
ACCENT_GOLD  = '#FFC72C'
ALERT_RED    = '#FF4D4D'
BORDER_GRAY  = '#252E42'

OUTDIR = '/Users/shuji/Desktop/kitt-voice/tmp_claude/slides_img_v3'
os.makedirs(OUTDIR, exist_ok=True)

def hex_to_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

BG_NAVY_RGB     = hex_to_rgb(BG_NAVY)
TEXT_WHITE_RGB  = hex_to_rgb(TEXT_WHITE)
TEXT_MUTED_RGB  = hex_to_rgb(TEXT_MUTED)
ACCENT_GOLD_RGB = hex_to_rgb(ACCENT_GOLD)
ALERT_RED_RGB   = hex_to_rgb(ALERT_RED)
BORDER_GRAY_RGB = hex_to_rgb(BORDER_GRAY)

JP_FONT = 'Hiragino Sans'

# 1920x1080 ⇒ Inches換算 (1 inch = 96 px)
# slide size: 16:9 → 13.333 x 7.5 inches
W_INCH = 13.333
H_INCH = 7.5
PX_PER_INCH = 1920 / W_INCH  # 144 px/inch

def px_to_in(px):
    return Inches(px / PX_PER_INCH)

prs = Presentation()
prs.slide_width  = Inches(W_INCH)
prs.slide_height = Inches(H_INCH)

# =================== ヘルパー ===================
def set_slide_bg(slide, color=BG_NAVY_RGB):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, x_px, y_px, w_px, h_px,
             size=20, color=TEXT_WHITE_RGB, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(px_to_in(x_px), px_to_in(y_px), px_to_in(w_px), px_to_in(h_px))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left  = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top   = Inches(0.05)
    tf.margin_bottom= Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = JP_FONT
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    return tb

def add_multiline(slide, lines, x_px, y_px, w_px, h_px, line_spacing=1.3):
    tb = slide.shapes.add_textbox(px_to_in(x_px), px_to_in(y_px), px_to_in(w_px), px_to_in(h_px))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top   = Inches(0.05)
    tf.margin_bottom= Inches(0.05)
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.line_spacing = line_spacing
        if isinstance(line, dict):
            text   = line.get('text', '')
            csize  = line.get('size', 20)
            ccolor = line.get('color', TEXT_WHITE_RGB)
            cbold  = line.get('bold', False)
        else:
            text, csize, ccolor, cbold = line, 20, TEXT_WHITE_RGB, False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.name = JP_FONT
        run.font.size = Pt(csize)
        run.font.color.rgb = ccolor
        run.font.bold = cbold

def add_card(slide, x_px, y_px, w_px, h_px, bg=BORDER_GRAY_RGB, line_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px_to_in(x_px), px_to_in(y_px), px_to_in(w_px), px_to_in(h_px))
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    if line_color is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = line_color
        card.line.width = Pt(1)
    card.shadow.inherit = False
    return card

def add_slide_footer(slide, page_no):
    add_text(slide, 'Project danjer-GAIA', 80, 1020, 800, 40, size=12, color=TEXT_MUTED_RGB)
    add_text(slide, f'{page_no:02d} / 09', 1700, 1020, 140, 40, size=12, color=TEXT_MUTED_RGB, align=PP_ALIGN.RIGHT)

def save_fig(fn, w=12, h=6.5, dpi=160, transparent=True):
    path = os.path.join(OUTDIR, fn)
    plt.gcf().set_size_inches(w, h)
    if not transparent:
        plt.gcf().patch.set_facecolor(BG_NAVY)
    else:
        plt.gcf().patch.set_alpha(0)
    plt.tight_layout(pad=0.5)
    plt.savefig(path, dpi=dpi, transparent=transparent, bbox_inches='tight')
    plt.close()
    return path

# =================== 図解 ===================
def fig_s2_paradigm():
    fig, ax = plt.subplots(figsize=(14, 6))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 6)
    ax.axis('off')

    # 上段: 固定bot
    ax.text(0.3, 5.3, '従来の固定bot', fontsize=14, color=TEXT_WHITE, weight='bold')
    x = np.linspace(1, 13, 100)
    y_old = 4.2 - 0.05 * x - 0.01 * (x ** 1.5) + 0.3 * np.sin(x * 0.7)
    ax.plot(x, y_old, color=ALERT_RED, linewidth=2.5, linestyle='--', alpha=0.85)
    ax.scatter([13], [y_old[-1]], s=200, color=ALERT_RED, zorder=5)
    ax.text(13.1, y_old[-1] - 0.1, 'ドローダウン破綻', color=ALERT_RED, fontsize=11, va='top')
    ax.text(7, 3.3, '相場の顔が変わった瞬間に無防備になる', ha='center', color=TEXT_MUTED, fontsize=11, style='italic')

    # 段間ラベル
    ax.text(7, 2.8, '─── レジーム識別による自動最適化 ───', ha='center', color=ACCENT_GOLD, fontsize=12, weight='bold')

    # 下段: GAIA
    ax.text(0.3, 2.3, '学び続けるGAIA', fontsize=14, color=ACCENT_GOLD, weight='bold')
    y_gaia = 1.0 + 0.12 * x + 0.2 * np.sin(x * 0.8)
    ax.plot(x, y_gaia, color=ACCENT_GOLD, linewidth=3, alpha=0.95)
    # 15分グリッド線
    for i in range(2, 14, 2):
        ax.axvline(i, ymin=0.02, ymax=0.35, color=TEXT_MUTED, linewidth=0.5, alpha=0.4)
    ax.scatter([13], [y_gaia[-1]], s=200, color=ACCENT_GOLD, zorder=5, edgecolor=TEXT_WHITE, linewidth=1.5)
    ax.text(13.1, y_gaia[-1] + 0.15, '右肩上がりに維持', color=ACCENT_GOLD, fontsize=11)

    return save_fig('s2_paradigm.png', w=14, h=6)


def fig_s3_funnel():
    fig, ax = plt.subplots(figsize=(8, 7))
    ax.set_xlim(-1, 1)
    ax.set_ylim(0, 8)
    ax.axis('off')

    # レイヤー1: 広口 (上)
    layer1 = Polygon([(-0.9, 6.5), (0.9, 6.5), (0.65, 5.5), (-0.65, 5.5)],
                     facecolor=TEXT_MUTED, alpha=0.35, edgecolor=BORDER_GRAY, linewidth=1.5)
    ax.add_patch(layer1)
    ax.text(0, 6.0, '全ポスト\n(雑談・感情・相場観)', ha='center', va='center',
            color=TEXT_WHITE, fontsize=11, weight='bold')

    # レイヤー2: 中口
    layer2 = Polygon([(-0.65, 5.4), (0.65, 5.4), (0.4, 4.4), (-0.4, 4.4)],
                     facecolor=BORDER_GRAY, alpha=0.85, edgecolor=TEXT_MUTED, linewidth=1.5)
    ax.add_patch(layer2)
    ax.text(0, 4.9, '投資判断ポスト\n失敗の反省', ha='center', va='center',
            color=TEXT_WHITE, fontsize=11, weight='bold')

    # レイヤー3: 最狭口 (下)
    layer3 = Polygon([(-0.4, 4.3), (0.4, 4.3), (0.25, 3.3), (-0.25, 3.3)],
                     facecolor=ACCENT_GOLD, alpha=0.95, edgecolor='#FFE082', linewidth=2)
    ax.add_patch(layer3)
    ax.text(0, 3.8, 'danjer DNA\n(確信度の源泉)', ha='center', va='center',
            color=BG_NAVY, fontsize=11, weight='bold')

    # 出力矢印
    ax.add_patch(FancyArrowPatch((0, 3.2), (0, 2.3), arrowstyle='->', mutation_scale=25,
                                  color=ACCENT_GOLD, linewidth=3))
    ax.text(0, 2.0, 'Slow Brain (頭脳) へ', ha='center', va='center',
            color=ACCENT_GOLD, fontsize=13, weight='bold')

    return save_fig('s3_funnel.png', w=8, h=7)


def fig_s4_radar():
    fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(projection=None))
    ax.set_xlim(-1.5, 1.5)
    ax.set_ylim(-1.5, 1.5)
    ax.axis('off')
    ax.set_aspect('equal')

    # 中心
    cx, cy = 0, 0
    r = 1.2

    # 3つの軸 (120度ずつ)
    axes_info = [
        (90, '方向\n(どこへ動くか)', ACCENT_GOLD),
        (90 - 120, '確信度\n(どれだけ強く張るか)', ACCENT_GOLD),
        (90 + 120, '危険度\n(最悪のシナリオ)', ALERT_RED),
    ]
    for angle_deg, label, col in axes_info:
        rad = np.deg2rad(angle_deg)
        x_end = cx + r * np.cos(rad)
        y_end = cy + r * np.sin(rad)
        ax.add_patch(FancyArrowPatch((cx, cy), (x_end, y_end), arrowstyle='-|>', mutation_scale=30,
                                      color=col, linewidth=3.5))
        label_r = r + 0.25
        ax.text(cx + label_r * np.cos(rad), cy + label_r * np.sin(rad), label,
                ha='center', va='center', color=col, fontsize=13, weight='bold')

    # 中央のサンプル現在値 (確信度・危険度・方向)
    sample = [0.7, 0.6, 0.3]  # 方向+0.7, 確信度+0.6, 危険度+0.3
    pts = []
    for (angle_deg, _, _), v in zip(axes_info, sample):
        rad = np.deg2rad(angle_deg)
        pts.append((cx + r * v * np.cos(rad), cy + r * v * np.sin(rad)))
    pts.append(pts[0])
    xs, ys = zip(*pts)
    ax.fill(xs, ys, color=ACCENT_GOLD, alpha=0.2)
    ax.plot(xs, ys, color=ACCENT_GOLD, linewidth=2)

    # 軸ラベル目盛
    for v in [0.25, 0.5, 0.75, 1.0]:
        for angle_deg, _, _ in axes_info:
            rad = np.deg2rad(angle_deg)
            ax.scatter([cx + r * v * np.cos(rad)], [cy + r * v * np.sin(rad)],
                       s=15, color=TEXT_MUTED, alpha=0.5)

    return save_fig('s4_radar.png', w=8, h=8)


def fig_s5_matrix():
    fig, ax = plt.subplots(figsize=(11, 7))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 7)
    ax.axis('off')

    # 2x2セル
    cells = [
        (1.0, 4.0, ACCENT_GOLD, '凪 × 上昇', '慎重ロング',  '小ロット + トレイリングTP'),
        (5.5, 4.0, ALERT_RED,   '凪 × 下降', '小ロット様子見', '反転の兆候待ち'),
        (1.0, 1.0, ACCENT_GOLD, '嵐 × 上昇', '一気に踏み込む', 'Half-Kelly上限フル'),
        (5.5, 1.0, ALERT_RED,   '嵐 × 下降', '退避またはショート', '高速SLで防衛'),
    ]
    for x, y, color, title, action, sub in cells:
        rect = Rectangle((x, y), 4.2, 2.5, facecolor=BORDER_GRAY, edgecolor=color, linewidth=2.2)
        ax.add_patch(rect)
        ax.text(x + 2.1, y + 2.05, title, ha='center', va='center', color=color, fontsize=14, weight='bold')
        ax.text(x + 2.1, y + 1.4, action, ha='center', va='center', color=TEXT_WHITE, fontsize=13, weight='bold')
        ax.text(x + 2.1, y + 0.5, sub, ha='center', va='center', color=TEXT_MUTED, fontsize=10)

    # 軸ラベル
    ax.text(0.5, 5.2, '凪 (低ボラ)', ha='center', va='center', color=TEXT_MUTED, fontsize=11, rotation=90)
    ax.text(0.5, 2.2, '嵐 (高ボラ)', ha='center', va='center', color=TEXT_MUTED, fontsize=11, rotation=90)
    ax.text(3.1, 0.5, '上昇トレンド', ha='center', va='center', color=TEXT_MUTED, fontsize=11)
    ax.text(7.6, 0.5, '下降トレンド', ha='center', va='center', color=TEXT_MUTED, fontsize=11)

    return save_fig('s5_matrix.png', w=11, h=7)


def fig_s6_triple_gate():
    fig, ax = plt.subplots(figsize=(14, 6))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 6)
    ax.axis('off')

    gates = [
        (0.5, '頭脳 (Slow Brain)', '15分間隔で\n戦略を練る', ACCENT_GOLD),
        (5.0, '安全装置 (Risk Engine)', '発注直前に\nレバ・損切幅を強制', TEXT_WHITE),
        (9.5, '反射神経 (Fast Guard)', '異常事態に\nミリ秒で自動ブレーキ', ALERT_RED),
    ]
    for x, title, body, color in gates:
        rect = FancyBboxPatch((x, 1.5), 3.5, 3.0, boxstyle='round,pad=0.05',
                              facecolor=BORDER_GRAY, edgecolor=color, linewidth=2.5)
        ax.add_patch(rect)
        ax.text(x + 1.75, 3.9, title, ha='center', va='center', color=color, fontsize=14, weight='bold')
        ax.text(x + 1.75, 3.0, body, ha='center', va='center', color=TEXT_WHITE, fontsize=11.5)
        # 役割タグ
        role_y = 2.0
        if 'Risk Engine' in title:
            ax.text(x + 1.75, role_y, '静的・構造的', ha='center', va='center', color=TEXT_MUTED, fontsize=10, style='italic')
        elif 'Fast Guard' in title:
            ax.text(x + 1.75, role_y, '動的・リアクティブ', ha='center', va='center', color=TEXT_MUTED, fontsize=10, style='italic')
        else:
            ax.text(x + 1.75, role_y, '思考・予測', ha='center', va='center', color=TEXT_MUTED, fontsize=10, style='italic')

    # 矢印
    for x in [4.0, 8.5]:
        ax.add_patch(FancyArrowPatch((x + 0.15, 3.0), (x + 0.85, 3.0), arrowstyle='->', mutation_scale=25,
                                      color=TEXT_WHITE, linewidth=2))

    # 発注矢印
    ax.add_patch(FancyArrowPatch((13.0, 3.0), (13.8, 3.0), arrowstyle='->', mutation_scale=25,
                                  color=ACCENT_GOLD, linewidth=2.5))

    # 必須メッセージ
    ax.text(7, 0.8, 'SLなしの注文は出さない — エントリーと同時に損切りを置く',
            ha='center', va='center', color=ACCENT_GOLD, fontsize=14, weight='bold')

    return save_fig('s6_gate.png', w=14, h=6)


def fig_s7_hourglass():
    fig, ax = plt.subplots(figsize=(6, 7))
    ax.set_xlim(-1.2, 1.2)
    ax.set_ylim(0, 8)
    ax.axis('off')

    # カウントダウンタイマー
    ax.text(0, 7.3, '15:00', ha='center', va='center', color=ACCENT_GOLD, fontsize=42, weight='bold')
    ax.text(0, 6.5, '判断の賞味期限', ha='center', va='center', color=TEXT_MUTED, fontsize=14)

    # 砂時計
    # 上部三角形 (砂が落ち始め、3割残っている設定)
    top_triangle = Polygon([(-0.8, 6.0), (0.8, 6.0), (0, 4.0)], facecolor=ACCENT_GOLD, alpha=0.7, edgecolor=ACCENT_GOLD)
    ax.add_patch(top_triangle)
    # 上部の砂 (残っている部分のみ薄く)
    sand_top = Polygon([(-0.55, 5.4), (0.55, 5.4), (0.1, 4.05), (-0.1, 4.05)], facecolor=ACCENT_GOLD, alpha=0.5)
    ax.add_patch(sand_top)

    # 中央の細い管 (くびれ)
    ax.plot([-0.05, -0.05], [4.0, 3.9], color=ACCENT_GOLD, linewidth=2)
    ax.plot([0.05, 0.05], [4.0, 3.9], color=ACCENT_GOLD, linewidth=2)

    # 下部三角形
    bot_triangle = Polygon([(0, 3.9), (-0.8, 1.9), (0.8, 1.9)], facecolor=BORDER_GRAY, alpha=0.5, edgecolor=ACCENT_GOLD)
    ax.add_patch(bot_triangle)
    # 下部に落ちた砂
    sand_bot = Polygon([(-0.55, 2.5), (0.55, 2.5), (0.7, 2.0), (-0.7, 2.0)], facecolor=ACCENT_GOLD, alpha=0.85)
    ax.add_patch(sand_bot)

    # 下にラベル
    ax.text(0, 1.5, '時間が来たら方針は失効', ha='center', va='center', color=TEXT_WHITE, fontsize=12)
    ax.text(0, 1.0, '応答停止5分で全閉', ha='center', va='center', color=ALERT_RED, fontsize=13, weight='bold')

    return save_fig('s7_hourglass.png', w=6, h=7)


def fig_s8_gantt():
    fig, ax = plt.subplots(figsize=(14, 6.5))
    ax.set_xlim(0, 14.5)
    ax.set_ylim(0, 7.5)
    ax.axis('off')

    # 7つの土台 (上から)
    tasks = [
        ('Trade-EHR 計算エンジン',          1,  3,  ACCENT_GOLD),
        ('過剰適合テスト (CPCV)',           1,  3,  ACCENT_GOLD),
        ('danjer DNA 抽出基盤',              4,  7,  TEXT_WHITE),
        ('レジーム判定 (相場の4つの顔)',     4,  7,  TEXT_WHITE),
        ('15分TTL+階層化スタンス',           8, 11,  ALERT_RED),
        ('注文前検問所',                     8, 11,  ALERT_RED),
        ('BQ棚卸し+朝サマリー',             12, 14,  ACCENT_GOLD),
    ]

    for i, (name, day_start, day_end, color) in enumerate(tasks):
        y = 6.5 - i * 0.85
        ax.text(0.1, y + 0.2, name, va='center', color=TEXT_WHITE, fontsize=12)
        # バー
        rect = Rectangle((day_start + 4.5, y), day_end - day_start + 1, 0.5,
                         facecolor=color, alpha=0.85, edgecolor=color)
        ax.add_patch(rect)

    # マイルストーン ピン
    milestones = [
        (3 + 4.5, 'Day 3: 評価基盤完成', ACCENT_GOLD),
        (7 + 4.5, 'Day 7: DNA+レジーム', TEXT_WHITE),
        (11 + 4.5, 'Day 11: 判断+防衛', ALERT_RED),
        (14 + 4.5, 'Day 14: 朝サマリー', ACCENT_GOLD),
    ]
    for x, label, color in milestones:
        ax.plot([x, x], [0.5, 7], color=color, linewidth=1.5, linestyle=':', alpha=0.6)
        ax.scatter([x], [0.6], s=100, color=color, marker='v', zorder=5)
        ax.text(x, 0.3, label, ha='center', va='top', color=color, fontsize=9.5)

    # Day軸 (横軸)
    for d in [1, 3, 7, 11, 14]:
        ax.text(d + 4.5, 7.2, f'D{d}', ha='center', color=TEXT_MUTED, fontsize=10)

    return save_fig('s8_gantt.png', w=14, h=6.5)


# =================== スライド本体 ===================
print('Generating slides...')

# ---------- S1: Project の核 ----------
print('S1...')
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
# コードネーム上部小
add_text(s, 'Project danjer-GAIA', 80, 80, 1200, 50,
         size=18, color=ACCENT_GOLD_RGB, bold=True)
# 大タイトル
add_text(s, 'BTC 投資効率最大化 AI', 80, 160, 1700, 130,
         size=60, color=TEXT_WHITE_RGB, bold=True)
# サブリード
add_text(s, '— danjerを教科書に、 AIが市場ゲームで成長する',
         80, 320, 1700, 80, size=28, color=ACCENT_GOLD_RGB)

# 3つのコアメッセージカード
core_msgs = [
    ('成長型', 'bot ではなく、 相場のルール変更に\n適応し続けるAI', TEXT_WHITE_RGB),
    ('右側予測', 'チャート未来を 「方向 × 確信度 × 危険度」\nの3軸で立体的に握る', ACCENT_GOLD_RGB),
    ('投資効率', '時間あたり利益 × 元手効率\n( = 時給 × 元手効率 = Trade-EHR )', ACCENT_GOLD_RGB),
]
for i, (head, sub, col) in enumerate(core_msgs):
    x = 80 + i * 620
    add_card(s, x, 530, 580, 280, bg=BORDER_GRAY_RGB, line_color=col)
    add_text(s, head, x + 30, 555, 540, 60, size=28, color=col, bold=True)
    add_text(s, sub, x + 30, 625, 540, 180, size=18, color=TEXT_WHITE_RGB)

# 下: タイムライン示唆
add_text(s, '8ヶ月で 検証可能な弟子AI まで進める / Day 1-14 で 動く土台 7つ',
         80, 870, 1760, 50, size=18, color=TEXT_MUTED_RGB)

add_slide_footer(s, 1)

# ---------- S2: 固定bot vs 学び続けるGAIA ----------
print('S2...')
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_text(s, '相場の変化に適応し続ける 「弟子AI」', 80, 80, 1760, 80,
         size=44, color=TEXT_WHITE_RGB, bold=True)
add_text(s, '固定ロジックbotは 賞味期限切れで破綻 — GAIAは 相場のルール変更に学び直す',
         80, 175, 1760, 50, size=20, color=TEXT_MUTED_RGB)

# 図解
fig_s2_paradigm()
s.shapes.add_picture(os.path.join(OUTDIR, 's2_paradigm.png'),
                      px_to_in(80), px_to_in(260), width=px_to_in(1760))

# 下: 補足
add_text(s, 'Shujiさんが嫌うのは 「成長しない固定bot」 — GAIAは自動実行を否定しない',
         80, 940, 1760, 60, size=18, color=ACCENT_GOLD_RGB, bold=True, align=PP_ALIGN.CENTER)
add_slide_footer(s, 2)

# ---------- S3: danjer DNAを「質」で絞り、その上を行く ----------
print('S3...')
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_text(s, 'danjer DNA の抽出プロセス', 80, 80, 1760, 80,
         size=44, color=TEXT_WHITE_RGB, bold=True)
add_text(s, '「教科書」を軽視しない。 だからこそ、 ノイズと戦略を徹底的に分離する',
         80, 175, 1760, 50, size=20, color=TEXT_MUTED_RGB)

# 漏斗図
fig_s3_funnel()
s.shapes.add_picture(os.path.join(OUTDIR, 's3_funnel.png'),
                      px_to_in(560), px_to_in(280), width=px_to_in(800))

# 左説明 (GPT Round 18 指摘反映: R27 明文化、 Shujiさんに通じる表現で)
add_multiline(s, [
    {'text': '「質」で絞る ≠ 売買発言だけ', 'size': 18, 'color': ACCENT_GOLD_RGB, 'bold': True},
    {'text': '常駐対象に含めるもの:', 'size': 17, 'color': TEXT_WHITE_RGB},
    {'text': '・売買判断ポスト', 'size': 16, 'color': TEXT_WHITE_RGB},
    {'text': '・失敗の反省', 'size': 16, 'color': TEXT_WHITE_RGB},
    {'text': '・警戒姿勢 (罠の匂い)', 'size': 16, 'color': TEXT_WHITE_RGB},
    {'text': '・相場観・利確判断', 'size': 16, 'color': TEXT_WHITE_RGB},
    {'text': '※ 雑談 / 個人事情は除外', 'size': 15, 'color': TEXT_MUTED_RGB},
], 80, 360, 460, 360)

# 右説明
add_multiline(s, [
    {'text': '完成したDNAは', 'size': 18, 'color': ACCENT_GOLD_RGB, 'bold': True},
    {'text': 'Slow Brain (頭脳) が', 'size': 18, 'color': TEXT_WHITE_RGB},
    {'text': '戦略を組み立てるための', 'size': 18, 'color': TEXT_WHITE_RGB},
    {'text': '絶対的なベースライン (地図)', 'size': 18, 'color': ACCENT_GOLD_RGB, 'bold': True},
    {'text': 'となる。 そしてAIはそれを超える', 'size': 18, 'color': TEXT_WHITE_RGB},
], 1380, 380, 480, 320)

add_slide_footer(s, 3)

# ---------- S4: チャートの右側を立体ホールド ----------
print('S4...')
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_text(s, '未来 (右側) を捉える 3つの座標軸', 80, 80, 1760, 80,
         size=44, color=TEXT_WHITE_RGB, bold=True)
add_text(s, '1つの未来を当てるのではなく、 方向×確信度×危険度の3軸で 「立体ホールド」 する',
         80, 175, 1760, 50, size=20, color=TEXT_MUTED_RGB)

fig_s4_radar()
s.shapes.add_picture(os.path.join(OUTDIR, 's4_radar.png'),
                      px_to_in(660), px_to_in(280), width=px_to_in(600))

# ゴールドメッセージ (S4 GPT必須)
add_card(s, 80, 870, 1760, 100, bg=BORDER_GRAY_RGB, line_color=ACCENT_GOLD_RGB)
add_text(s, '確信度が 高く、 かつ 危険度が 低い 局面 (圧倒的優位性) においてのみ、 レバレッジを 動的に引き上げる',
         80, 870, 1760, 100, size=22, color=ACCENT_GOLD_RGB, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_slide_footer(s, 4)

# ---------- S5: 相場の4つの顔 ----------
print('S5...')
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_text(s, 'レジーム (相場の4つの顔) とマトリクス', 80, 80, 1760, 80,
         size=44, color=TEXT_WHITE_RGB, bold=True)
add_text(s, 'ボラ × トレンドの4象限で レバを 自動で 上下する — 攻めるべき時だけ攻める',
         80, 175, 1760, 50, size=20, color=TEXT_MUTED_RGB)

fig_s5_matrix()
s.shapes.add_picture(os.path.join(OUTDIR, 's5_matrix.png'),
                      px_to_in(280), px_to_in(280), width=px_to_in(1360))

add_text(s, 'レジーム判定は ATR × Slope を ローリング・パーセンタイル で 動的に閾値設定',
         80, 950, 1760, 50, size=16, color=TEXT_MUTED_RGB, align=PP_ALIGN.CENTER, bold=False)
add_slide_footer(s, 5)

# ---------- S6: 三重の防護壁 ----------
print('S6...')
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_text(s, '損切りは AIの判断より前 — 三重の防護壁', 80, 80, 1760, 80,
         size=42, color=TEXT_WHITE_RGB, bold=True)
add_text(s, 'AIが判断する前に 3つの関門で 「破滅」 を 物理的に止める',
         80, 175, 1760, 50, size=20, color=TEXT_MUTED_RGB)

fig_s6_triple_gate()
s.shapes.add_picture(os.path.join(OUTDIR, 's6_gate.png'),
                      px_to_in(80), px_to_in(280), width=px_to_in(1760))

# 必須メッセージ重ね (GPT指示)
add_card(s, 80, 870, 1760, 100, bg=BORDER_GRAY_RGB, line_color=ACCENT_GOLD_RGB)
add_text(s, 'SL なしの注文は 出さない。 エントリーと同時に 損切りを 置く',
         80, 870, 1760, 100, size=24, color=ACCENT_GOLD_RGB, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_slide_footer(s, 6)

# ---------- S7: 15分の賞味期限 ----------
print('S7...')
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_text(s, '15分 TTL (賞味期限) と デッドロック解消', 80, 80, 1760, 80,
         size=42, color=TEXT_WHITE_RGB, bold=True)
add_text(s, '古い判断で 勝手に売買しない — Slow Brainの方針には 鉄の時間制限がある',
         80, 175, 1760, 50, size=20, color=TEXT_MUTED_RGB)

fig_s7_hourglass()
s.shapes.add_picture(os.path.join(OUTDIR, 's7_hourglass.png'),
                      px_to_in(80), px_to_in(290), width=px_to_in(600))

# 右側 ロジック
add_text(s, '時間による 鉄の統治', 760, 290, 1080, 60,
         size=32, color=ACCENT_GOLD_RGB, bold=True)
add_multiline(s, [
    {'text': '判断の賞味期限は ', 'size': 22, 'color': TEXT_WHITE_RGB},
    {'text': '「15分」', 'size': 28, 'color': ACCENT_GOLD_RGB, 'bold': True},
    {'text': '— TTL切れ後は 新規エントリーを 完全禁止', 'size': 18, 'color': TEXT_WHITE_RGB},
], 760, 380, 1080, 130)

add_multiline(s, [
    {'text': 'Slow Brainの応答停止が ', 'size': 22, 'color': TEXT_WHITE_RGB},
    {'text': '「5分」', 'size': 28, 'color': ALERT_RED_RGB, 'bold': True},
    {'text': '続いた場合、 異常事態とみなし安全装置作動', 'size': 18, 'color': TEXT_WHITE_RGB},
    {'text': '→ 全ポジションの 縮小・強制クローズ', 'size': 22, 'color': ALERT_RED_RGB, 'bold': True},
], 760, 540, 1080, 200)

add_text(s, '※ 一括成行による APIロック (429エラー) を回避するため、 ミリ秒単位のキューイング制御',
         760, 770, 1080, 50, size=14, color=ALERT_RED_RGB)

add_slide_footer(s, 7)

# ---------- S8: Day 1-14 ガントチャート ----------
print('S8...')
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_text(s, 'Day 1 - 14 開発・実装スケジュール', 80, 80, 1760, 80,
         size=42, color=TEXT_WHITE_RGB, bold=True)
add_text(s, '巨大設計書ではなく 「動くもの」 を 14日間で 7つ 作る',
         80, 175, 1760, 50, size=20, color=TEXT_MUTED_RGB)

fig_s8_gantt()
s.shapes.add_picture(os.path.join(OUTDIR, 's8_gantt.png'),
                      px_to_in(80), px_to_in(260), width=px_to_in(1760))

add_text(s, '※ Trade-EHR から着手 (評価軸が無いと他が無意味)',
         80, 950, 1760, 50, size=16, color=ACCENT_GOLD_RGB, align=PP_ALIGN.CENTER)
add_slide_footer(s, 8)

# ---------- S9: Shujiさんへの3大決断 ----------
print('S9...')
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_text(s, 'Shujiさんに 本日決断いただく 3つの 大枠方向性', 80, 80, 1760, 80,
         size=38, color=TEXT_WHITE_RGB, bold=True)
add_text(s, '実弾金額や細かな数値は Phase 3 以降に持ち越し — まず方向と優先順位だけ',
         80, 175, 1760, 50, size=20, color=TEXT_MUTED_RGB)

decisions = [
    ('01',
     'コードネーム',
     'Project danjer-GAIA で 進めるか',
     'GPT/Gemini/Claude 3者推奨',
     ACCENT_GOLD_RGB),
    ('02',
     '最優先フェーズ',
     'danjer DNA を 「質ベース」 で 抽出し、 Day 1 から着手してよいか',
     '※ 件数で絞らず 質で絞る (Shujiさん Q1 回答反映)',
     ACCENT_GOLD_RGB),
    ('03',
     '運用順序',
     '紙トレ → 0.0001 BTC 小額実弾 → 完全自動売買 の順 で進めるか',
     '※ いきなり実弾はしない',
     ACCENT_GOLD_RGB),
]
for i, (num, head, body, note, col) in enumerate(decisions):
    y = 290 + i * 220
    add_card(s, 80, y, 1760, 200, bg=BORDER_GRAY_RGB, line_color=col)
    # 番号
    add_text(s, num, 120, y + 30, 130, 130,
             size=44, color=col, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # heading
    add_text(s, head, 280, y + 25, 1500, 50,
             size=22, color=col, bold=True)
    # body
    add_text(s, body, 280, y + 80, 1500, 60,
             size=24, color=TEXT_WHITE_RGB, bold=True)
    # note
    add_text(s, note, 280, y + 145, 1500, 50,
             size=16, color=TEXT_MUTED_RGB)

add_slide_footer(s, 9)


# ============================================================
# Appendix Slides (S10-S12): 専門用語フル解禁版 + 右側に用語辞書
# ============================================================

# 共通: 用語辞書サイドバーを描画するヘルパー
def add_glossary_sidebar(slide, glossary_items, x_px=1320, y_px=130, w_px=520, h_px=870):
    """右側 用語辞書"""
    add_card(slide, x_px, y_px, w_px, h_px, bg=BORDER_GRAY_RGB, line_color=ACCENT_GOLD_RGB)
    # ヘッダ
    add_text(slide, '用語辞書', x_px + 20, y_px + 15, w_px - 40, 45,
             size=18, color=ACCENT_GOLD_RGB, bold=True)
    # 項目を整形
    lines = []
    for term, desc in glossary_items:
        lines.append({'text': term, 'size': 13, 'color': ACCENT_GOLD_RGB, 'bold': True})
        lines.append({'text': desc, 'size': 11, 'color': TEXT_WHITE_RGB})
    add_multiline(slide, lines, x_px + 20, y_px + 65, w_px - 40, h_px - 80, line_spacing=1.2)


# ---------- S10: v1/v2/v3 戦略比較表 ----------
print('S10...')
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_text(s, 'v1 / v2 / v3 戦略比較 — 何が変わったか', 80, 80, 1760, 60,
         size=32, color=TEXT_WHITE_RGB, bold=True)

# 比較表 (左コンテンツ領域 80-1280 px)
# テーブル: 行=戦略軸、列=v1/v2/v3
table_data = [
    ('軸',           'v1 (14ラウンド)',       'v2 (徹夜)',                'v3 (本気3者)'),
    ('議論方式',     '14ラウンド',            '7ラウンド',                '20ラウンド ぐるぐる無制限'),
    ('投資効率指標', 'SAC-Lagrangian + 月コスト目標', 'ProfitPerHour (バグ含)',      'Trade-EHR (分母ガード+MA30)'),
    ('danjer移植',   'DNA-0〜3段階, BC試験',  'RAG+常駐 49,667件全件',    '質ベース絞り込み (投資判断のみ常駐)'),
    ('アーキ',       '10エージェント + Particle Filter', '3つの魂 (攻め/守り/師匠)',     'GAIA-Triad 2.0 (4ノード)'),
    ('右側予測',     'Scenario DSL 述語25', 'Probabilistic Deep Hedging', '3軸 (方向/確信度/危険度)'),
    ('レバ',         'Half-Kelly',           'Half-Kelly + 制約',         'Half-Kelly × Confidence × Vol'),
    ('リスク',       'Risk Engine 6層',       '注文前検問 6step',           'Fast Guard / Risk Engine 二層'),
    ('デッドロック', '未検討',                '未検討',                    'TTL + 階層化スタンス (15分)'),
    ('検証',         'CPCV N=10 / DSR / PSR / PBO', 'CPCV + EVT 併用案',          'CPCV + 動的閾値 + EVT'),
    ('多取引所',     'Bybit + Hyperliquid + OKX', '同',                        '同 + Rate Limit 対策'),
    ('月コスト',     '$135-400 (Phase別)',    '$50-100 (Cache想定)',       '$25-200 (POC後確定)'),
    ('Phase',        'Phase 1-5+ 8ヶ月',     '同',                       'Day 1-14 動く土台7つ → Phase以降'),
    ('議事録',       'round_table.md',       'round_table_v2.md',        'round_table_v3.md (Live)'),
    ('Shuji評価',    '10年前の新人レベル',    'プロ品質、 ただし2人会議',   '本気3者、 動線改善、 質ベース反映'),
]

# テーブル幅: 1200 / 各列: 軸250 + v1 310 + v2 310 + v3 320
COL_W = [250, 310, 310, 330]
COL_X = [80]
for w in COL_W[:-1]:
    COL_X.append(COL_X[-1] + w)

# ヘッダ
header_y = 165
for i, (header, w) in enumerate(zip(table_data[0], COL_W)):
    add_card(s, COL_X[i], header_y, w - 4, 50, bg=BORDER_GRAY_RGB, line_color=ACCENT_GOLD_RGB)
    add_text(s, header, COL_X[i] + 10, header_y + 5, w - 24, 40,
             size=13, color=ACCENT_GOLD_RGB, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 行
for ri, row in enumerate(table_data[1:]):
    y = header_y + 60 + ri * 50
    # 行 stripe
    for ci, (cell, w) in enumerate(zip(row, COL_W)):
        bg_col = BORDER_GRAY_RGB if ri % 2 == 0 else hex_to_rgb('#1a2238')
        line_col = None
        if ci == 0:
            text_col = ACCENT_GOLD_RGB if ri == len(table_data) - 2 else TEXT_WHITE_RGB
            font_size = 11
            font_bold = True
        else:
            text_col = ALERT_RED_RGB if (ri == len(table_data) - 2 and ci == 1) else (ACCENT_GOLD_RGB if (ri == len(table_data) - 2 and ci >= 2) else TEXT_WHITE_RGB)
            font_size = 10
            font_bold = (ri == len(table_data) - 2)
        add_card(s, COL_X[ci], y, w - 4, 45, bg=bg_col, line_color=None)
        add_text(s, cell, COL_X[ci] + 10, y + 4, w - 24, 38,
                 size=font_size, color=text_col, bold=font_bold, anchor=MSO_ANCHOR.MIDDLE)

# 右: 用語辞書
glossary_s10 = [
    ('SAC-Lagrangian', '制約付きSoft Actor-Critic。リスク制約を満たしつつ報酬最大化する強化学習。'),
    ('Trade-EHR', 'Equity-Hours Return = NetProfit/(AvgEquity×ElapsedHours)。 時給×元手効率。'),
    ('MA30', '直近30トレードの移動平均。EHRの揺らぎを平準化。'),
    ('RAG', 'Retrieval-Augmented Generation。 検索拡張生成。 過去ポストを類似検索。'),
    ('GAIA-Triad 2.0', 'Slow Brain+Risk Engine+Fast Guard+Order Gate の4ノード。'),
    ('Particle Filter', '粒子フィルタ。 複数シナリオを並走させ観測で確率を更新。'),
    ('Half-Kelly', 'ケリー基準の半分。 安全側にレバを抑える。 p-(1-p)/b の半分。'),
    ('TTL', 'Time-To-Live。 判断の賞味期限。 v3 では15分。'),
    ('CPCV', 'Combinatorial Purged Cross-Validation。 過学習検出BT。'),
    ('DSR/PSR/PBO', 'Deflated/Probabilistic Sharpe Ratio / Probability of Backtest Overfitting。 偶然勝ち検出。'),
    ('EVT', 'Extreme Value Theory。 極値統計。 ファットテール対策。'),
    ('Probabilistic Deep Hedging', '複数シナリオに確率分布と最適ヘッジ比率を同時最適化。'),
]
add_glossary_sidebar(s, glossary_s10, x_px=1320, y_px=165, w_px=540, h_px=870)
add_slide_footer(s, 10)


# ---------- S11: v3 仕様書 (システム全体像) ----------
print('S11...')
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_text(s, 'v3 仕様書 — システム全体像 / 担当 / 使用ツール', 80, 80, 1760, 60,
         size=30, color=TEXT_WHITE_RGB, bold=True)

# 左コンテンツ領域 (80-1280)
# システム全体図 (matplotlib で描画)
def fig_s11_arch():
    fig, ax = plt.subplots(figsize=(13, 7))
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 9)
    ax.axis('off')

    # データ層 (下)
    ax.add_patch(FancyBboxPatch((0.5, 0.5), 12, 1.3, boxstyle='round,pad=0.05',
                                 facecolor=BORDER_GRAY, edgecolor=TEXT_MUTED, linewidth=1.5))
    ax.text(6.5, 1.5, 'データ層', ha='center', va='center', color=TEXT_MUTED, fontsize=12, weight='bold')
    ax.text(2.5, 0.9, 'BigQuery btc_trading', ha='center', va='center', color=TEXT_WHITE, fontsize=11, weight='bold')
    ax.text(2.5, 0.6, 'OHLCV/OI/FR/清算/LS/danjer DNA', ha='center', va='center', color=TEXT_MUTED, fontsize=9)
    ax.text(6.5, 0.9, 'Cloud Storage', ha='center', va='center', color=TEXT_WHITE, fontsize=11, weight='bold')
    ax.text(6.5, 0.6, '画像/ログ 14日', ha='center', va='center', color=TEXT_MUTED, fontsize=9)
    ax.text(10.5, 0.9, 'Bybit / Hyperliquid API', ha='center', va='center', color=TEXT_WHITE, fontsize=11, weight='bold')
    ax.text(10.5, 0.6, '価格・板・注文', ha='center', va='center', color=TEXT_MUTED, fontsize=9)

    # 推論層 (中央)
    boxes = [
        (0.5, 'Slow Brain',    'Gemini 3.1 Pro\nContext Cache (15分)', ACCENT_GOLD),
        (3.4, 'Risk Engine',   '静的検問\nEVT/MaxDD/PF',                TEXT_WHITE),
        (6.3, 'Fast Guard',    'ルールベース\nms単位ブレーキ',          ALERT_RED),
        (9.2, 'Order Gate',    '6ステップ検問\nSL/レバ/流動性',          ACCENT_GOLD),
    ]
    for x, name, body, col in boxes:
        ax.add_patch(FancyBboxPatch((x, 3.5), 2.8, 2.0, boxstyle='round,pad=0.05',
                                     facecolor=BG_NAVY, edgecolor=col, linewidth=2))
        ax.text(x + 1.4, 4.9, name, ha='center', va='center', color=col, fontsize=13, weight='bold')
        ax.text(x + 1.4, 4.0, body, ha='center', va='center', color=TEXT_WHITE, fontsize=10)

    # 矢印 (Slow Brain → Risk Engine → Fast Guard → Order Gate → 取引所)
    for x_start in [3.3, 6.2, 9.1]:
        ax.add_patch(FancyArrowPatch((x_start, 4.5), (x_start + 0.1, 4.5),
                                      arrowstyle='->', mutation_scale=20, color=TEXT_WHITE, linewidth=2))
    # 下から上へ (データ → 推論)
    ax.add_patch(FancyArrowPatch((6.5, 1.85), (6.5, 3.45), arrowstyle='->', mutation_scale=20, color=TEXT_MUTED, linewidth=1.5))
    # 上から下へ (Order Gate → 取引所)
    ax.add_patch(FancyArrowPatch((10.6, 3.45), (10.6, 1.85), arrowstyle='->', mutation_scale=20, color=ACCENT_GOLD, linewidth=2))

    # オペレータ層 (上)
    ax.add_patch(FancyBboxPatch((0.5, 6.5), 12, 1.8, boxstyle='round,pad=0.05',
                                 facecolor=BORDER_GRAY, edgecolor=TEXT_MUTED, linewidth=1.5))
    ax.text(6.5, 8.0, 'オペレータ層', ha='center', va='center', color=TEXT_MUTED, fontsize=12, weight='bold')
    ops = [
        (1.5, 7.0, 'Shujiさん',    '最終判断・観察・仕様変更', ACCENT_GOLD),
        (4.5, 7.0, 'Claude Code', '実装・BQ・デプロイ・議事録', TEXT_WHITE),
        (7.5, 7.0, 'GPT (Web)',   '司会・戦略統合', TEXT_WHITE),
        (10.5, 7.0, 'Gemini (Web)', '監査・常駐LLM・ビジュアル', ACCENT_GOLD),
    ]
    for x, y, name, role, col in ops:
        ax.text(x, y, name, ha='center', va='center', color=col, fontsize=12, weight='bold')
        ax.text(x, y - 0.4, role, ha='center', va='center', color=TEXT_WHITE, fontsize=9)

    return save_fig('s11_arch.png', w=13, h=7)

fig_s11_arch()
s.shapes.add_picture(os.path.join(OUTDIR, 's11_arch.png'),
                      px_to_in(80), px_to_in(160), width=px_to_in(1200))

# 下: 担当×ツール表
add_text(s, '担当 × 使用ツール', 80, 760, 1200, 40,
         size=16, color=ACCENT_GOLD_RGB, bold=True)
tool_table = [
    ('Shujiさん', 'iPhone Safari / Keynote / GitHub', '最終判断 / 仕様変更指示 / pptx確認'),
    ('Claude Code', 'Bash / Python / python-pptx / BigQuery / gcloud / git', '実装全般 / BQ DDL / Cloud Run デプロイ / 議事録'),
    ('GPT (ChatGPT Web)', 'ChatGPT (gemini-2.0-flashとは別)', '司会 / 戦略統合 / 整合性監査'),
    ('Gemini (Gemini Web + 3.1 Pro Context Cache)', 'Gemini Web (常駐LLM) / Imagen 3', '監査 / Slow Brain本番運用 / ビジュアル設計'),
]
for i, (who, tool, role) in enumerate(tool_table):
    y = 805 + i * 50
    add_text(s, who, 80, y, 350, 40, size=11, color=ACCENT_GOLD_RGB, bold=True)
    add_text(s, tool, 430, y, 400, 40, size=10, color=TEXT_WHITE_RGB)
    add_text(s, role, 830, y, 450, 40, size=10, color=TEXT_MUTED_RGB)

# 右: 用語辞書
glossary_s11 = [
    ('BigQuery (BQ)', 'Google大規模データウェアハウス。1TiB free tier。'),
    ('Context Cache', 'Gemini API のトークン常駐機能。 通常入力の1/4料金。'),
    ('Cloud Run', 'GCPのステートレスコンテナ実行サービス。 e2-microは無料tier。'),
    ('Slow Brain', '15分間隔で戦略 (スタンスJSON) を出力する 「頭脳」 ノード。 Gemini 3.1 Pro Context Cache。'),
    ('Risk Engine', '発注前に静的ルールで危険注文を止める検問所。'),
    ('Fast Guard', 'ms単位で異常を検知し、 ブレーキのみ独自判断するルールベース層。'),
    ('Order Gate', '注文前 6ステップ最終検問。 SL/レバ/流動性/コスト/類似Pattern/Explainability。'),
    ('OHLCV', 'Open/High/Low/Close/Volume。 ローソク足の基本データ。'),
    ('OI', 'Open Interest。 未決済建玉。'),
    ('FR', 'Funding Rate。 永久先物の資金調達率。'),
    ('LS比', 'Long/Short ratio。 ポジション偏り指標。'),
    ('Imagen 3', 'Google画像生成AI (Geminiパイプライン)。'),
]
add_glossary_sidebar(s, glossary_s11, x_px=1320, y_px=160, w_px=540, h_px=870)
add_slide_footer(s, 11)


# ---------- S12: コスト見積り明細 ----------
print('S12...')
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_text(s, 'コスト見積り明細 — 固定費 / 変動費', 80, 80, 1760, 60,
         size=32, color=TEXT_WHITE_RGB, bold=True)
add_text(s, 'Phase 1-2 POC段階 / Phase 3-5 拡大時 で 2フェーズ提示', 80, 150, 1200, 40,
         size=14, color=TEXT_MUTED_RGB)

# 固定費テーブル (Phase 1-2)
add_text(s, '【固定費】(毎月発生、 利用量に依らず)', 80, 200, 600, 40,
         size=16, color=ACCENT_GOLD_RGB, bold=True)
fixed = [
    ('GitHub',          '無料',        '議事録/コード/pptx ホスト'),
    ('BigQuery',        '$0 (1TiB free)', 'クエリ1TiB/月まで無料'),
    ('Cloud Run (e2-micro)', '$0-5',   '常時稼働分。 free tier 内なら$0'),
    ('Cloud Storage',   '$0-2',        '画像/ログ 14日ライフサイクル、 ~10GB想定'),
    ('Cloud Scheduler', '$0',          '5ジョブまで無料'),
]
header_y2 = 245
add_card(s, 80, header_y2, 1200, 36, bg=BORDER_GRAY_RGB, line_color=ACCENT_GOLD_RGB)
add_text(s, 'サービス', 90, header_y2 + 4, 300, 28, size=12, color=ACCENT_GOLD_RGB, bold=True)
add_text(s, '月額',    400, header_y2 + 4, 200, 28, size=12, color=ACCENT_GOLD_RGB, bold=True)
add_text(s, '内訳',    610, header_y2 + 4, 660, 28, size=12, color=ACCENT_GOLD_RGB, bold=True)
for i, (svc, cost, body) in enumerate(fixed):
    y = header_y2 + 40 + i * 36
    add_card(s, 80, y, 1200, 32, bg=hex_to_rgb('#1a2238') if i % 2 else BORDER_GRAY_RGB, line_color=None)
    add_text(s, svc, 90, y + 2, 300, 28, size=11, color=TEXT_WHITE_RGB, bold=True)
    add_text(s, cost, 400, y + 2, 200, 28, size=11, color=ACCENT_GOLD_RGB)
    add_text(s, body, 610, y + 2, 660, 28, size=10, color=TEXT_MUTED_RGB)
add_text(s, '固定費合計 (Phase 1-2): 月 $0-7',
         80, header_y2 + 40 + len(fixed) * 36 + 8, 1200, 32,
         size=14, color=ACCENT_GOLD_RGB, bold=True)

# 変動費テーブル
var_y = 480
add_text(s, '【変動費】(利用量に応じて、 取引活動・モデル学習量に比例)', 80, var_y, 1200, 40,
         size=16, color=ACCENT_GOLD_RGB, bold=True)
variable = [
    ('Gemini 3.1 Pro Context Cache 常駐', '$10-20',      '49,667ポストの30-40% (~300K tokens) を24/7常駐'),
    ('Gemini クエリ (15分間隔)',       '$5-30',          '日96クエリ × 30日 = 2,880回/月、 出力tokens変動'),
    ('Gemini Embedding API (DNA初期化)', '$2-5 (一回限り)', '49,667ポスト × 1024dim embedding。 Phase 1で1回'),
    ('GPU Spot (A100 1h/週)',           '$6-15',          '週1回 RL/PBT 学習。 Spot価格変動あり'),
    ('Bybit/Hyperliquid API',           '$0',             '読取無料、 取引手数料は別'),
    ('取引手数料 (Phase 5以降)',          '別途',           'Maker -0.025% / Taker 0.075% × 取引量'),
]
add_card(s, 80, var_y + 45, 1200, 36, bg=BORDER_GRAY_RGB, line_color=ACCENT_GOLD_RGB)
add_text(s, 'サービス', 90, var_y + 49, 380, 28, size=12, color=ACCENT_GOLD_RGB, bold=True)
add_text(s, '月額',    480, var_y + 49, 180, 28, size=12, color=ACCENT_GOLD_RGB, bold=True)
add_text(s, '内訳',    670, var_y + 49, 600, 28, size=12, color=ACCENT_GOLD_RGB, bold=True)
for i, (svc, cost, body) in enumerate(variable):
    y = var_y + 85 + i * 36
    add_card(s, 80, y, 1200, 32, bg=hex_to_rgb('#1a2238') if i % 2 else BORDER_GRAY_RGB, line_color=None)
    add_text(s, svc, 90, y + 2, 380, 28, size=10, color=TEXT_WHITE_RGB, bold=True)
    add_text(s, cost, 480, y + 2, 180, 28, size=10, color=ACCENT_GOLD_RGB)
    add_text(s, body, 670, y + 2, 600, 28, size=9, color=TEXT_MUTED_RGB)

add_text(s, '変動費合計 (Phase 1-2): 月 $25-65  /  Phase 3-5 拡大時: 月 $50-200',
         80, var_y + 85 + len(variable) * 36 + 8, 1200, 32,
         size=14, color=ACCENT_GOLD_RGB, bold=True)

# 合計バナー
add_card(s, 80, 950, 1200, 70, bg=BORDER_GRAY_RGB, line_color=ACCENT_GOLD_RGB)
add_text(s, '【総額】 Phase 1-2 POC: 月 $25-72  /  Phase 3-5 本格運用: 月 $50-207',
         80, 950, 1200, 70, size=18, color=ACCENT_GOLD_RGB, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 右: 用語辞書
glossary_s12 = [
    ('POC', 'Proof of Concept。 概念実証。 実コストや実機能を最小実装で測る段階。'),
    ('Free Tier', 'GCPの無料利用枠。 BigQuery 1TiB/月、 Cloud Run e2-micro 等。'),
    ('Context Cache', 'Gemini入力tokenを24時間キャッシュ。 通常料金の1/4。'),
    ('Embedding API', 'テキスト/画像をベクトル化するAPI。 類似検索の前処理。'),
    ('GPU Spot', '中断可能だが安価なGPUインスタンス。 A100 で 約 $1.5/h。'),
    ('A100', 'NVIDIA データセンターGPU。 RL学習で主流。'),
    ('PBT', 'Population Based Training。 個体並走進化学習。'),
    ('Maker / Taker', '指値約定 (Maker、 流動性提供) / 成行約定 (Taker、 流動性消費) の手数料区分。'),
    ('Funding Rate', '永久先物の資金調達料。 ロングが多いとロングがショートに支払う。'),
    ('OHLCV取得', 'ローソク足データ。 Bybit/Hyperliquid REST APIで無料取得可。'),
    ('Cloud Scheduler', 'GCPのcron代替。 5ジョブまで無料。'),
    ('Cloud Run', 'コンテナ実行サービス。 リクエスト時のみ起動で課金最小化。'),
]
add_glossary_sidebar(s, glossary_s12, x_px=1320, y_px=160, w_px=540, h_px=870)
add_slide_footer(s, 12)


# 保存
OUT = '/Users/shuji/Desktop/kitt-voice/btc_ai_master_plan_v3.pptx'
prs.save(OUT)
print(f'OK: {OUT}')
